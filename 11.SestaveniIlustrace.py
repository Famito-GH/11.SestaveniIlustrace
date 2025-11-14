import pandas as pd
from pptx import Presentation
import os
import sys
import comtypes.client
import pythoncom
import logging
import glob
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import tkinter as tk
from tkinter import messagebox, Listbox, MULTIPLE, END
import threading

# ---------------- LOGGING ---------------- #
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

log_file = os.path.join(BASE_DIR, "log.txt")

logging.basicConfig(
    filename=log_file,
    level=logging.DEBUG,
    format="%(asctime)s - %(levelname)s - %(message)s",
    filemode="w"
)

logging.info("===== Spouštím skript 11.SestaveniIlustrace ===== " + BASE_DIR)

# ---------------- SOUBORY ---------------- #
excel_files = glob.glob(os.path.join(BASE_DIR, "*.xlsx"))
pptx_files = glob.glob(os.path.join(BASE_DIR, "*.pptx"))

excel_file = excel_files[0] if excel_files else None
pptx_file = pptx_files[0] if pptx_files else None

if not excel_file:
    logging.error("❌ Nebyl nalezen žádný Excel soubor (.xlsx) ve složce: " + BASE_DIR)
if not pptx_file:
    logging.error("❌ Nebyl nalezen žádný PowerPoint soubor (.pptx) ve složce: " + BASE_DIR)

output_folder = os.path.join(BASE_DIR, "exported_slides")
os.makedirs(output_folder, exist_ok=True)

# ---------------- NAČTENÍ EXCELU ---------------- #
df = pd.DataFrame()
if excel_file:
    try:
        df = pd.read_excel(excel_file, header=0)  # hlavička na prvním řádku
        df.columns = [str(c).strip() for c in df.columns]
        logging.info(f"Excel úspěšně načten. Sloupce: {list(df.columns-2)}")
    except Exception as e:
        logging.exception(f"❌ Chyba při načítání Excelu: {e}")

def format_excel_value(val):
    if pd.isna(val):
        return ""
    # Pokud je číslo
    if isinstance(val, (int, float)):
        # Pokud je float a je celé, zobraz jako int
        if isinstance(val, float) and val.is_integer():
            return str(int(val))
        return str(val)
    return str(val).strip()

# ---------------- EXPORT ---------------- #
slides_processed = 0
export_in_progress = False

def export_selected_products(selected_kody=None):
    global slides_processed, export_in_progress
    
    # Inicializuj COM pro tento thread
    pythoncom.CoInitialize()
    
    try:
        slides_processed = 0
        export_in_progress = True

        base_dir = os.path.dirname(sys.executable) if getattr(sys, 'frozen', False) else os.path.dirname(os.path.abspath(__file__))
        excel_files = glob.glob(os.path.join(base_dir, "*.xlsx"))
        pptx_files = glob.glob(os.path.join(base_dir, "*.pptx"))

        excel_file = excel_files[0] if excel_files else None
        pptx_file = pptx_files[0] if pptx_files else None

        if not excel_file or not pptx_file:
            logging.error("❌ Nebyl nalezen Excel nebo PowerPoint soubor.")
            messagebox.showerror("Chyba", "Nebyl nalezen Excel nebo PowerPoint soubor ve složce.")
            export_in_progress = False
            return

        try:
            df = pd.read_excel(excel_file, header=0)
            df.columns = [str(c).strip() for c in df.columns]
            df["Číslo modelu"] = df["Číslo modelu"].astype(str).str.strip()
            if "Kód" in df.columns:
                df["Kód"] = df["Kód"].astype(str).str.strip()
        except Exception as e:
            logging.exception(f"❌ Chyba při načítání Excelu: {e}")
            messagebox.showerror("Chyba", f"Chyba při načítání Excelu: {e}")
            export_in_progress = False
            return

        required_columns = ["Číslo modelu", "Hmotnost (kg)", "ŠÍŘKA", "VÝŠKA", "HLOUBKA", 
                           "Šířka popruhu", "Maximální délka popruhu", "Minimální délka popruhu"]
        dropna_columns = [col for col in required_columns if col != "Číslo modelu"]
        if "Kód" in df.columns:
            dropna_columns.append("Kód")
        valid_rows = df.dropna(subset=dropna_columns)

        if selected_kody and "Kód" in df.columns:
            selected_kody_str = [str(k).strip() for k in selected_kody]
            valid_rows = valid_rows[valid_rows["Kód"].isin(selected_kody_str)]

        if valid_rows.empty:
            logging.info("ℹ️ Žádné validní řádky k exportu.")
            messagebox.showinfo("Hotovo", "Nebyl nalezen žádný produkt k exportu.")
            export_in_progress = False
            return

        shape_to_column = {
            "váha": "Hmotnost (kg)", "šířka": "ŠÍŘKA", "výška": "VÝŠKA", "hloubka": "HLOUBKA",
            "šířka popruhu": "Šířka popruhu", "max. délka popruhu": "Maximální délka popruhu",
            "min. délka popruhu": "Minimální délka popruhu", "objem": "Objem",
            "výška ucha": "Výška ucha", "ramenní popruhy": "Šířka ucha", "šířka ucha": "Šířka ucha",
            "šířka uch": "Šířka ucha", "ucho základna": "Ucho základna", "CisloModelu": "Číslo modelu",
        }

        grouped = valid_rows.groupby("Číslo modelu")
        total_to_process = sum(len(group) for _, group in grouped)
        logging.info(f"Začínám export: {total_to_process} produktů ve {len(grouped)} skupinách podle čísla modelu.")

        powerpoint = None
        try:
            # Inicializuj PowerPoint jednou pro všechny exporty
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            powerpoint.WindowState = 2  # 2 = minimalizované okno
            
            for model, group in grouped:
                model_str = str(model).strip()
                logging.info(f"--- Zpracovávám model: {model_str} (položek: {len(group)}) ---")

                for _, row in group.iterrows():
                    kod = str(row["Kód"]).strip() if "Kód" in row and not pd.isna(row["Kód"]) else "NEZNAMY"
                    presentation = None
                    try:
                        prs = Presentation(pptx_file)
                        target_slide_idx = None
                        
                        for i, slide in enumerate(prs.slides):
                            for shape in slide.shapes:
                                if getattr(shape, "name", "") == "CisloModelu":
                                    slide_model = str(getattr(shape, "text", "")).strip()
                                    try:
                                        model_num = float(model_str)
                                        model_str_comp = str(int(model_num)) if model_num.is_integer() else str(model_num)
                                    except:
                                        model_str_comp = str(model_str)
                                    
                                    if slide_model == model_str_comp:
                                        target_slide_idx = i
                                        break
                            if target_slide_idx is not None:
                                break

                        if target_slide_idx is None:
                            logging.warning(f"⚠️ Šablonový slide pro model '{model_str}' nebyl nalezen. Kód {kod} přeskočen.")
                            continue

                        target_slide = prs.slides[target_slide_idx]

                        for shape in target_slide.shapes:
                            shape_name = getattr(shape, "name", "")
                            if shape_name in shape_to_column:
                                excel_col = shape_to_column[shape_name]
                                value = format_excel_value(row.get(excel_col, ""))
                                if value == "" or value.lower() == "nan":
                                    continue

                                if shape_name == "váha":
                                    value_str = f"{value} kg"
                                elif shape_name in ["šířka", "výška", "hloubka", "šířka popruhu",
                                                    "max. délka popruhu", "min. délka popruhu", 
                                                    "ramenní popruhy", "výška ucha", "šířka ucha", "ucho základna"]:
                                    value_str = f"{value} cm"
                                elif shape_name == "objem":
                                    value_str = f"{value} l"
                                elif shape_name == "CisloModelu":
                                    value_str = model_str
                                else:
                                    value_str = value

                                try:
                                    if hasattr(shape, "text"):
                                        shape.text = value_str
                                    elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                                        shape.text_frame.clear()
                                        shape.text_frame.paragraphs[0].text = value_str
                                except Exception as e:
                                    logging.debug(f"Nelze nastavit text pro shape '{shape_name}': {e}")

                                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                                    if shape_name in ["šířka popruhu", "hloubka", "váha", "min. délka popruhu", "objem", 
                                                      "ramenní popruhy","max. délka popruhu", "výška ucha", "šířka ucha", "ucho základna"]:
                                        for paragraph in shape.text_frame.paragraphs:
                                            paragraph.alignment = PP_ALIGN.RIGHT
                                    for paragraph in shape.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.name = "Open Sans"
                                            run.font.bold = True
                                            if shape_name in ["šířka", "výška", "hloubka", "šířka popruhu",
                                                              "max. délka popruhu", "min. délka popruhu", "ramenní popruhy", 
                                                              "objem","výška ucha", "šířka ucha", "ucho základna"]:
                                                run.font.size = Pt(44)
                                            elif shape_name == "váha":
                                                run.font.size = Pt(28)

                        temp_pptx = os.path.join(output_folder, f"__temp_{kod}.pptx")
                        prs.save(temp_pptx)

                        try:
                            presentation = powerpoint.Presentations.Open(temp_pptx, WithWindow=False)
                            export_path = os.path.join(output_folder, f"{kod}_20.jpg")
                            presentation.Slides.Item(target_slide_idx + 1).Export(export_path, "JPG")
                            logging.info(f"✅ Kód {kod}: exportováno do {export_path}")
                            slides_processed += 1
                        except Exception as export_err:
                            logging.error(f"Chyba při exportu {kod}: {export_err}")
                        finally:
                            if presentation:
                                try:
                                    presentation.Close(0)
                                except:
                                    pass
                                presentation = None

                    except Exception as e:
                        logging.exception(f"❌ Chyba při zpracování kódu {kod} (model {model_str}): {e}")
                    finally:
                        if presentation:
                            try:
                                presentation.Close(0)
                            except:
                                pass

        finally:
            if powerpoint:
                try:
                    powerpoint.Quit()
                except:
                    pass
                powerpoint = None

        # Smaž všechny __temp_* soubory až po uzavření PowerPointu
        logging.info("Odstraňuji dočasné soubory...")
        temp_files = glob.glob(os.path.join(output_folder, "__temp_*"))
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
                logging.info(f"✓ Smazán: {temp_file}")
            except Exception as e:
                logging.warning(f"⚠️ Nelze smazat {temp_file}: {e}")

        logging.info(f"Zpracováno {slides_processed} slidů.")
        logging.info("===== Konec skriptu =====")
        
        # Vynutit garbage collection
        import gc
        gc.collect()
        
        export_in_progress = False
        messagebox.showinfo("Hotovo", f"Zpracováno {slides_processed} slidů.\nVýstup: {output_folder}")
    finally:
        # Vždy vyčisti COM
        pythoncom.CoUninitialize()
        export_in_progress = False

# ---------------- GUI ---------------- #
def gui_main():
    kody = []

    def update_products():
        nonlocal kody
        listbox.delete(0, END)
        if excel_file:
            try:
                df_gui = pd.read_excel(excel_file, header=0)
                df_gui.columns = [str(c).strip() for c in df_gui.columns]
                if "Kód" in df_gui.columns:
                    kody = df_gui["Kód"].dropna().astype(str).unique().tolist()
                    for k in kody:
                        listbox.insert(END, k)
            except:
                pass

    root = tk.Tk()
    root.title("Sestavení Ilustrace - Export produktů")
    root.geometry("500x470")

    mode_var = tk.IntVar(value=0)

    frame_mode = tk.Frame(root)
    frame_mode.pack(anchor="w", padx=10, pady=(10,0))

    tk.Label(frame_mode, text="Vyberte režim exportu:").pack(anchor="w")
    tk.Radiobutton(frame_mode, text="Všechny produkty", variable=mode_var, value=0, command=lambda: toggle_listbox()).pack(anchor="w")
    tk.Radiobutton(frame_mode, text="Konkrétní produkty", variable=mode_var, value=1, command=lambda: toggle_listbox()).pack(anchor="w")

    listbox_label = tk.Label(root, text="Vyberte produkty:")
    listbox = Listbox(root, selectmode=MULTIPLE, height=15)
    for k in kody:
        listbox.insert(END, k)

    def toggle_listbox():
        if mode_var.get() == 0:
            listbox_label.pack_forget()
            listbox.pack_forget()
            btn_update.pack_forget()
        else:
            listbox_label.pack(anchor="w", padx=10, pady=(10,0))
            listbox.pack(fill="both", expand=True, padx=10, pady=(0,10))
            btn_update.pack(fill="x", padx=10, pady=(0,10))

    btn_update = tk.Button(root, text="Aktualizovat produkty", command=update_products)
    btn_export = tk.Button(root, text="Exportovat", bg="#4CAF50", fg="white", height=2)
    toggle_listbox()

    def run_export_threaded():
        global export_in_progress
        if export_in_progress:
            messagebox.showwarning("Upozornění", "Export již probíhá!")
            return
        
        btn_export.config(state="disabled", text="Vytvářejí se ilustrace...")
        
        if mode_var.get() == 0:
            thread = threading.Thread(target=export_selected_products, args=(None,), daemon=False)
        else:
            selected = [listbox.get(i) for i in listbox.curselection()]
            if not selected:
                messagebox.showerror("Chyba", "Vyberte alespoň jeden produkt.")
                btn_export.config(state="normal", text="Exportovat")
                return
            thread = threading.Thread(target=export_selected_products, args=(selected,), daemon=False)
        
        thread.start()
        
        # Periodicky kontroluj, zda export skončil
        def check_export():
            if export_in_progress:
                root.after(500, check_export)
            else:
                btn_export.config(state="normal", text="Exportovat")
        
        root.after(500, check_export)

    btn_export.config(command=run_export_threaded)
    btn_export.pack(fill="x", padx=10, pady=10)

    root.mainloop()


if __name__ == "__main__":
    gui_main()
    sys.exit()
